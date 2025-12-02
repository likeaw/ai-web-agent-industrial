"""
Microsoft Office 文档操作工具。

支持创建和操作以下文档类型：
- Word: .docx, .doc
- Excel: .xlsx, .xls
- PowerPoint: .pptx, .ppt

注意：.doc, .xls, .ppt 是旧格式，本工具主要支持新格式（.docx, .xlsx, .pptx）。
对于旧格式，会尝试创建新格式文件。
"""

import os
from typing import List, Tuple, Optional, Dict, Any
from pathlib import Path

# 导入 Office 文档处理库（如果未安装会抛出 ImportError）
try:
    from docx import Document
    from docx.shared import Pt, RGBColor
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from .file_operations import check_path_safety, resolve_user_path


def _check_office_library(file_type: str) -> Tuple[bool, Optional[str]]:
    """检查所需的 Office 库是否已安装。"""
    if file_type in ["docx", "doc"]:
        if not DOCX_AVAILABLE:
            return False, "python-docx library is not installed. Install it with: pip install python-docx"
    elif file_type in ["xlsx", "xls"]:
        if not EXCEL_AVAILABLE:
            return False, "openpyxl library is not installed. Install it with: pip install openpyxl"
    elif file_type in ["pptx", "ppt"]:
        if not PPTX_AVAILABLE:
            return False, "python-pptx library is not installed. Install it with: pip install python-pptx"
    return True, None


def create_word_document(path: str, content: Optional[str] = None, title: Optional[str] = None) -> Tuple[bool, str]:
    """
    创建 Word 文档（.docx）。

    :param path: 文档保存路径（应包含 .docx 扩展名）。
    :param content: 文档内容（纯文本，支持多行）。
    :param title: 文档标题（可选）。
    :return: (是否成功, 结果消息)
    """
    lib_ok, lib_error = _check_office_library("docx")
    if not lib_ok:
        return False, lib_error

    try:
        abs_path = resolve_user_path(path)
    except ValueError as exc:
        return False, f"Invalid path: {exc}"

    is_safe, error_msg = check_path_safety(abs_path, "write")
    if not is_safe:
        return False, f"Safety check failed: {error_msg}"

    try:
        # 确保扩展名是 .docx
        if not abs_path.lower().endswith(".docx"):
            if abs_path.lower().endswith(".doc"):
                abs_path = abs_path[:-4] + ".docx"
            else:
                abs_path += ".docx"

        # 确保父目录存在
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)

        # 创建文档
        doc = Document()

        # 添加标题（如果有）
        if title:
            title_para = doc.add_heading(title, level=1)
            title_para.alignment = 1  # 居中

        # 添加内容（如果有）
        if content:
            # 按行分割内容
            lines = content.split("\n")
            for line in lines:
                if line.strip():
                    para = doc.add_paragraph(line.strip())
                else:
                    doc.add_paragraph()  # 空行

        # 保存文档
        doc.save(abs_path)
        return True, f"Word document created: {abs_path}"

    except Exception as e:
        return False, f"Failed to create Word document: {e}"


def create_excel_document(
    path: str,
    data: Optional[List[List[Any]]] = None,
    sheet_name: str = "Sheet1",
    headers: Optional[List[str]] = None,
) -> Tuple[bool, str]:
    """
    创建 Excel 文档（.xlsx）。

    :param path: 文档保存路径（应包含 .xlsx 扩展名）。
    :param data: 数据行列表，每行是一个列表，例如：[["A1", "B1"], ["A2", "B2"]]。
    :param sheet_name: 工作表名称（默认 "Sheet1"）。
    :param headers: 表头列表（可选），例如：["姓名", "年龄"]。
    :return: (是否成功, 结果消息)
    """
    lib_ok, lib_error = _check_office_library("xlsx")
    if not lib_ok:
        return False, lib_error

    try:
        abs_path = resolve_user_path(path)
    except ValueError as exc:
        return False, f"Invalid path: {exc}"

    is_safe, error_msg = check_path_safety(abs_path, "write")
    if not is_safe:
        return False, f"Safety check failed: {error_msg}"

    try:
        # 确保扩展名是 .xlsx
        if not abs_path.lower().endswith(".xlsx"):
            if abs_path.lower().endswith(".xls"):
                abs_path = abs_path[:-4] + ".xlsx"
            else:
                abs_path += ".xlsx"

        # 确保父目录存在
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)

        # 创建工作簿
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name

        # 添加表头（如果有）
        if headers:
            for col_idx, header in enumerate(headers, start=1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal="center", vertical="center")

        # 添加数据（如果有）
        start_row = 2 if headers else 1
        if data:
            for row_idx, row_data in enumerate(data, start=start_row):
                for col_idx, cell_value in enumerate(row_data, start=1):
                    ws.cell(row=row_idx, column=col_idx, value=cell_value)

        # 保存工作簿
        wb.save(abs_path)
        return True, f"Excel document created: {abs_path}"

    except Exception as e:
        return False, f"Failed to create Excel document: {e}"


def create_powerpoint_document(
    path: str,
    slides: Optional[List[Dict[str, Any]]] = None,
    title: Optional[str] = None,
) -> Tuple[bool, str]:
    """
    创建 PowerPoint 演示文稿（.pptx）。

    :param path: 文档保存路径（应包含 .pptx 扩展名）。
    :param slides: 幻灯片列表，每个幻灯片是一个字典，包含：
                   - "title": 幻灯片标题（可选）
                   - "content": 幻灯片内容（文本列表，每项是一段）
    :param title: 演示文稿标题（可选，用于第一张幻灯片）。
    :return: (是否成功, 结果消息)
    """
    lib_ok, lib_error = _check_office_library("pptx")
    if not lib_ok:
        return False, lib_error

    try:
        abs_path = resolve_user_path(path)
    except ValueError as exc:
        return False, f"Invalid path: {exc}"

    is_safe, error_msg = check_path_safety(abs_path, "write")
    if not is_safe:
        return False, f"Safety check failed: {error_msg}"

    try:
        # 确保扩展名是 .pptx
        if not abs_path.lower().endswith(".pptx"):
            if abs_path.lower().endswith(".ppt"):
                abs_path = abs_path[:-4] + ".pptx"
            else:
                abs_path += ".pptx"

        # 确保父目录存在
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)

        # 创建演示文稿
        prs = Presentation()

        # 如果没有提供幻灯片，至少创建一张标题幻灯片
        if not slides:
            if title:
                slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片布局
                title_shape = slide.shapes.title
                title_shape.text = title
            else:
                # 创建空白幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        else:
            # 添加每张幻灯片
            for slide_data in slides:
                # 使用标题和内容布局
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                content_shape = slide.placeholders[1]

                # 设置标题
                slide_title = slide_data.get("title", "")
                if not slide_title and title and slides.index(slide_data) == 0:
                    slide_title = title
                title_shape.text = slide_title

                # 设置内容
                content = slide_data.get("content", [])
                if isinstance(content, str):
                    content = [content]
                text_frame = content_shape.text_frame
                text_frame.text = content[0] if content else ""
                for para_text in content[1:]:
                    p = text_frame.add_paragraph()
                    p.text = para_text

        # 保存演示文稿
        prs.save(abs_path)
        return True, f"PowerPoint document created: {abs_path}"

    except Exception as e:
        return False, f"Failed to create PowerPoint document: {e}"


def create_office_document(
    file_type: str,
    path: str,
    **kwargs
) -> Tuple[bool, str]:
    """
    通用 Office 文档创建函数，根据文件类型自动选择相应的创建函数。

    :param file_type: 文档类型（"docx", "xlsx", "pptx"）。
    :param path: 文档保存路径。
    :param kwargs: 其他参数，传递给相应的创建函数。
    :return: (是否成功, 结果消息)
    """
    file_type_lower = file_type.lower().lstrip(".")

    if file_type_lower in ["docx", "doc"]:
        return create_word_document(
            path,
            content=kwargs.get("content"),
            title=kwargs.get("title"),
        )
    elif file_type_lower in ["xlsx", "xls"]:
        return create_excel_document(
            path,
            data=kwargs.get("data"),
            sheet_name=kwargs.get("sheet_name", "Sheet1"),
            headers=kwargs.get("headers"),
        )
    elif file_type_lower in ["pptx", "ppt"]:
        return create_powerpoint_document(
            path,
            slides=kwargs.get("slides"),
            title=kwargs.get("title"),
        )
    else:
        return False, f"Unsupported file type: {file_type}. Supported types: docx, xlsx, pptx"

